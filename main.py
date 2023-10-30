import os
from pptx import Presentation


def delete_speaker_notes(presentation):
    for slide in presentation.slides:
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = ''


def process_powerpoint_files():
    current_directory = os.getcwd()
    modified_folder = os.path.join(current_directory, "Modified")

    # Create the "Modified" folder if it doesn't exist
    if not os.path.exists(modified_folder):
        os.makedirs(modified_folder)

    for filename in os.listdir(current_directory):
        if filename.endswith(".pptx") or filename.endswith(".ppt"):
            presentation_path = os.path.join(current_directory, filename)

            # Open the presentation
            presentation = Presentation(presentation_path)

            # Delete speaker notes from every slide
            delete_speaker_notes(presentation)

            # Save the modified presentation in the "Modified" folder
            modified_path = os.path.join(modified_folder, filename)
            presentation.save(modified_path)

            print(f"Speaker notes deleted from {filename}. Modified presentation saved as {modified_path}")


process_powerpoint_files()
